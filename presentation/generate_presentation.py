from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Action Spotting in SoccerNet-v2"
subtitle.text = "Architecture Optimization, Temporal Scaling, and Pipeline Automation\nPhD Supervisory Update"

# Slide 2: Motivation & Problem Statement
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "1. Motivation & Problem Statement"
tf = body.text_frame
p = tf.add_paragraph()
p.text = "Research Domain: Spatio-Temporal Action Spotting in broadcast football video."
p.level = 0
p = tf.add_paragraph()
p.text = "Core Challenge 1: Severe Class Imbalance."
p.level = 1
p = tf.add_paragraph()
p.text = "Background events and high-frequency actions (Throw-ins) completely drown out rare, critical events (Penalties, Red Cards)."
p.level = 2
p = tf.add_paragraph()
p.text = "Core Challenge 2: Temporal Complexity."
p.level = 1
p = tf.add_paragraph()
p.text = "Unlike static image classification, identifying a 'Clearance' or 'Shot' requires contextual mathematical memory across multiple frames."
p.level = 2

# Slide 3: Key Innovations
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "2. Key Innovations in Current Pipeline"
tf = body.text_frame
p = tf.add_paragraph()
p.text = "Algorithmic Shift: Replaced standard Cross-Entropy with Focal Loss (Gamma=2.0)."
p.level = 0
p = tf.add_paragraph()
p.text = "Mathematically penalizes the network for ignoring rare classes, fixing the model's tendency to predict a single dominant class."
p.level = 1
p = tf.add_paragraph()
p.text = "Architectural Strategy: Progressive Temporal Unfreezing."
p.level = 0
p = tf.add_paragraph()
p.text = "Strategically unfreezing deeper 3D CNN blocks to force the pre-trained SlowFast network to learn complex football-specific temporal mechanics."
p.level = 1

# Slide 4: Methodology
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "3. Methodology"
tf = body.text_frame
p = tf.add_paragraph()
p.text = "Architecture: SlowFast R50 (Dual-pathway 3D CNN for video)."
p.level = 0
p = tf.add_paragraph()
p.text = "Data Pipeline: Automated frame extraction and tensor standardization mapped to PyTorch Dataloaders."
p.level = 0
p = tf.add_paragraph()
p.text = "Evaluation Metric: Official SoccerNet Evaluation API to guarantee standardized Average-mAP (1s-5s tolerances)."
p.level = 0
p = tf.add_paragraph()
p.text = "Telemetry: Integrated Weights & Biases (W&B) for rigorous gradient flow and validation tracking."
p.level = 0

# Slide 5: Experimental Results
slide_layout = prs.slide_layouts[5] # Title only layout for table
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "4. Experimental Results"

x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(2)
table = slide.shapes.add_table(rows=4, cols=3, left=x, top=y, width=cx, height=cy).table
table.cell(0, 0).text = "Model Configuration"
table.cell(0, 1).text = "Official mAP"
table.cell(0, 2).text = "First Half mAP"

table.cell(1, 0).text = "Baseline (Frozen Backbone)"
table.cell(1, 1).text = "1.34%"
table.cell(1, 2).text = "1.88%"

table.cell(2, 0).text = "Exp 1: Unfrozen Block 5"
table.cell(2, 1).text = "1.79%"
table.cell(2, 2).text = "2.18%"

table.cell(3, 0).text = "Exp 2: Unfrozen Blocks 4 & 5"
table.cell(3, 1).text = "2.76%"
table.cell(3, 2).text = "4.49%"

# Add a text box below table
txBox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf2 = txBox.text_frame
p = tf2.add_paragraph()
p.text = "Finding: Doubling the network's unfrozen temporal capacity (Blocks 4 & 5) directly doubled the resulting mAP score."
p.level = 0

# Slide 6: Per-Class Breakthroughs
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "5. Per-Class Analysis (Exp 2)"

x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4)
table2 = slide.shapes.add_table(rows=6, cols=2, left=x, top=y, width=cx, height=cy).table
table2.cell(0, 0).text = "Action Class"
table2.cell(0, 1).text = "mAP Achieved"

table2.cell(1, 0).text = "Ball out of play"
table2.cell(1, 1).text = "17.68%"

table2.cell(2, 0).text = "Kick-off"
table2.cell(2, 1).text = "15.91% (New discovery)"

table2.cell(3, 0).text = "Corner"
table2.cell(3, 1).text = "14.08%"

table2.cell(4, 0).text = "Throw-in"
table2.cell(4, 1).text = "11.19%"

table2.cell(5, 0).text = "Clearance, Foul, Shots"
table2.cell(5, 1).text = "1% - 7% (New discoveries)"

# Add text below
txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(7), Inches(1))
tf3 = txBox2.text_frame
p = tf3.add_paragraph()
p.text = "Finding: Increased temporal capacity allowed the model to expand from recognizing 4 basic classes to 9 complex classes."
p.level = 0

# Slide 7: Discussion & Limitations
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "6. Discussion & Limitations"
tf = body.text_frame
p = tf.add_paragraph()
p.text = "Dataset Saturation (Overfitting):"
p.level = 0
p = tf.add_paragraph()
p.text = "In Exp 2, Training Loss dropped to 0.79 while Validation Loss rose to 1.98."
p.level = 1
p = tf.add_paragraph()
p.text = "The network has exhausted the variance of the local 5-game dataset and is memorizing it. Unfreezing further layers without more data will severely degrade generalization."
p.level = 1
p = tf.add_paragraph()
p.text = "Hardware Constraints:"
p.level = 0
p = tf.add_paragraph()
p.text = "Unfreezing Blocks 4 & 5 required drastically reducing the batch size to 2 to prevent GPU CUDA crashes, which negatively impacts gradient stability."
p.level = 1

# Slide 8: Future Work
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "7. Future Work"
tf = body.text_frame
p = tf.add_paragraph()
p.text = "Goal: Achieve SOTA double-digit mAP scores across all 17 target classes."
p.level = 0
p = tf.add_paragraph()
p.text = "1. Data Scaling:"
p.level = 1
p = tf.add_paragraph()
p.text = "Expand the training set from 5 games to the full 500-game SoccerNet-v2 dataset to solve the overfitting bottleneck."
p.level = 2
p = tf.add_paragraph()
p.text = "2. Infrastructure Migration:"
p.level = 1
p = tf.add_paragraph()
p.text = "Package and deploy the refactored, research-ready pipeline to the University HPC cluster."
p.level = 2
p = tf.add_paragraph()
p.text = "3. Batch Size Stabilization:"
p.level = 1
p = tf.add_paragraph()
p.text = "Utilize High-VRAM GPUs (24GB+) to increase the batch size to 16, allowing the entirely unfrozen architecture to be trained with stable gradients."
p.level = 2

prs.save("Supervisory_Meeting_Presentation_v3.pptx")
print("Academic Presentation generated successfully!")
